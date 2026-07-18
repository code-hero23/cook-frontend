import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize presentation & widescreen (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 2. Design System Colors
DARK_BG = RGBColor(15, 23, 42)        # Slate 900 (#0f172a)
CARD_BG = RGBColor(30, 41, 59)        # Slate 800 (#1e293b)
CODE_BG = RGBColor(2, 6, 23)          # Slate 950 (#020617)
TEXT_WHITE = RGBColor(248, 250, 252)  # Slate 50 (#f8fafc)
TEXT_MUTED = RGBColor(148, 163, 184)  # Slate 400 (#94a3b8)
ACCENT_ORANGE = RGBColor(249, 115, 22) # Orange 500 (#f97316)
BORDER_COLOR = RGBColor(51, 65, 85)    # Slate 700 (#334155)

# Helper function to clear layouts and apply full screen background
def create_slide_with_bg():
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Draw background rectangle
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK_BG
    rect.line.fill.background()
    return slide

# Helper function to add headers
def add_slide_header(slide, title_text, category_text=None):
    if category_text:
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Segoe UI"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_ORANGE
        
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

# Helper function to add a card layout component
def add_card(slide, left, top, width, height, title, body_points, bg_color=CARD_BG):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_color
    rect.line.color.rgb = BORDER_COLOR
    rect.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), height - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = ACCENT_ORANGE
    p_title.space_after = Pt(14)
    
    for pt in body_points:
        p = tf.add_paragraph()
        p.text = "• " + pt
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(8)

# ----------------- SLIDE 1: Title Slide -----------------
slide1 = create_slide_with_bg()

# Highlight design accent line
accent_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(4.5), Inches(0.08))
accent_line.fill.solid()
accent_line.fill.fore_color.rgb = ACCENT_ORANGE
accent_line.line.fill.background()

# Title text frame
title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.5))
tf1 = title_box.text_frame
tf1.word_wrap = True
tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0

p_title = tf1.paragraphs[0]
p_title.text = "UNIFIED COOKSCAPE"
p_title.font.name = "Segoe UI"
p_title.font.size = Pt(54)
p_title.font.bold = True
p_title.font.color.rgb = TEXT_WHITE
p_title.space_after = Pt(8)

p_sub = tf1.add_paragraph()
p_sub.text = "Enterprise-Grade Full-Stack Business Management, HRMS & CRM Platform"
p_sub.font.name = "Segoe UI"
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = TEXT_MUTED
p_sub.space_after = Pt(36)

p_tech = tf1.add_paragraph()
p_tech.text = "React 19  |  Tailwind CSS 4.0  |  Node.js (Express)  |  PostgreSQL & Prisma"
p_tech.font.name = "Segoe UI"
p_tech.font.size = Pt(13)
p_tech.font.bold = True
p_tech.font.color.rgb = ACCENT_ORANGE

# ----------------- SLIDE 2: Technical Architecture -----------------
slide2 = create_slide_with_bg()
add_slide_header(slide2, "Technical Architecture & System Stack", "SYSTEM OVERVIEW")

# 3 Columns for layers
add_card(slide2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Frontend Layer", [
    "Built with React 19 (Vite SPA) for fast compile & rendering speed.",
    "Tailwind CSS 4.0 utility-first framework for premium, modern styles.",
    "React Router 7 handles client-side views & permission auth guards.",
    "Framer Motion executes micro-animations & transitions.",
    "Recharts engine displays real-time operational dashboard analytics."
])

add_card(slide2, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Backend REST API", [
    "Node.js & Express.js server utilizing Controller-Service-Route structure.",
    "JSON Web Token (JWT) + Google OAuth 2.0 secure authorization.",
    "Multer middleware manages project assets and BGV documents.",
    "Node-cron automates background task checking and daily email reports.",
    "Web-Push API delivers instant real-time browser push notifications."
])

add_card(slide2, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), "Database & Integrations", [
    "Prisma ORM provides type-safe queries, transaction safety, and migration histories.",
    "PostgreSQL relational DB handles complex entities (Users, Tasks, Vouchers, Logs).",
    "Nodemailer handles automated operational notifications (SMTP Gmail).",
    "Google APIs interface for third-party OAuth and cloud Drive backups."
])

# ----------------- SLIDE 3: Admin & Manager Functionalities -----------------
slide3 = create_slide_with_bg()
add_slide_header(slide3, "Admin & Manager: Corporate Operations", "USER ROLES & FUNCTIONALITIES")

add_card(slide3, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "HRMS & BGV", [
    "Full Employee Lifecycle tracking (onboarding, documents, assets, exit).",
    "Background Verification (BGV) workflow with digital document upload.",
    "CSV employee import/export supporting automated parsing.",
    "Secure permission roles: SUPER_ADMIN, ADMIN, MANAGER, and EMPLOYEE."
])

add_card(slide3, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Voucher Approvals", [
    "Hierarchical expense voucher approval matrix safeguarding finances.",
    "Sequential lifecycle states:\nSubmitted → AM Review → COO Review → Paid → Completed.",
    "Audited activity log logs every stage transitions with timestamp details.",
    "Allows downloading summaries to Excel via ExcelJS integration."
])

add_card(slide3, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), "System Schedulers", [
    "Node-cron runs automated database backups sent directly to Google Drive.",
    "System checks daily for overdue tasks and sends warning emails.",
    "Automatically aggregates employee daily work reports and emails managers.",
    "Keep-alive triggers keep external API services awake and healthy."
])

# ----------------- SLIDE 4: CRE (Client Relations Executive) -----------------
slide4 = create_slide_with_bg()
add_slide_header(slide4, "CRE: Front-Office CRM & Showroom Hub", "USER ROLES & FUNCTIONALITIES")

add_card(slide4, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Walkin Hub CRM", [
    "Enables CREs to register walk-in client details upon showroom arrival.",
    "Tracks client details, showroom visited, tentative timings, and remarks.",
    "Automated WhatsApp message triggers notify Business Heads of new leads.",
    "Integrated status states to flag active, pending, or converted walk-ins."
])

add_card(slide4, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Daily Work Reports", [
    "CREs submit daily logs detailing client interactions and progress.",
    "Tracks next follow-up date (nextFD) to prompt active customer retention.",
    "Includes a client satisfaction rating system (1-5 stars).",
    "Automatically references Business Heads (BH) and Field Agents (FA)."
])

add_card(slide4, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), "Monthly Analytics", [
    "Aggregates total customer calls, site visits (SRV), and proposal counts.",
    "Computes signed customer order counts and revenue value generated.",
    "Provides performance data displayed in the Admin dashboard."
])

# ----------------- SLIDE 5: Supervisor & Employee -----------------
slide5 = create_slide_with_bg()
add_slide_header(slide5, "Supervisor & Employee: Field Task Execution", "USER ROLES & FUNCTIONALITIES")

add_card(slide5, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Task Execution & Stages", [
    "Assigned projects feature sequential task stages and due dates.",
    "Real-time task tracking updates task status: PENDING, IN_PROGRESS, COMPLETED.",
    "Auto-escalation of overdue tasks triggers supervisor warning emails."
])

add_card(slide5, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Geo-Tagged Evidence", [
    "Requires employees to upload photos proving on-site work completion.",
    "Captures and records device GPS location (latitude and longitude) on upload.",
    "Prevents spoofing by cross-referencing task evidence coordinates with project locations."
])

add_card(slide5, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), "Collaboration & Issues", [
    "In-project chat messaging facilitates communication between team and client.",
    "Ticketing system allows flagging blocking issues directly on tasks.",
    "Enables attaching PDF blueprints and images to comments."
])

# ----------------- SLIDE 6: Client Dashboard & Experience -----------------
slide6 = create_slide_with_bg()
add_slide_header(slide6, "Client Portal: Transparency & Feedback", "USER ROLES & FUNCTIONALITIES")

add_card(slide6, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Magic Link Authentication", [
    "Passwordless, secure token-based logins for end clients.",
    "Expiration-based links ensure client security without account friction.",
    "Allows instant access to personal project details from any device."
])

add_card(slide6, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Visual Progress Tracking", [
    "Real-time progress bars mapping actual project execution percentage.",
    "Transparent billing status displaying payments received vs. outstanding values.",
    "Shows completed project milestones with verified evidence photos."
])

add_card(slide6, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), "Interactive Channels", [
    "Allows clients to send messages directly to project supervisors.",
    "Direct ticket submission interface for complaints, additions, or changes.",
    "Provides real-time feedback with instant notifications."
])

# ----------------- SLIDE 7: Data Flow & Execution -----------------
slide7 = create_slide_with_bg()
add_slide_header(slide7, "Data Flow & System Execution", "TECHNICAL WORKFLOW")

# We will create a horizontal flow using cards or process blocks
flow_points = [
    ("1. Client Request", "React frontend sends request (e.g. Upload Task Evidence) via Axios with JWT header."),
    ("2. JWT Middleware", "Express backend decodes & validates the JWT. Verifies user roles and access rights."),
    ("3. Business Logic", "Controller processes files/JSON. Captures geo-location (Lat/Long) from uploaded metadata."),
    ("4. Database Update", "Prisma client saves data to PostgreSQL, triggering audit activity logs & db hooks."),
    ("5. Push & Mail Alert", "Nodemailer fires SMTP email alerts, Web-Push triggers client/admin notifications.")
]

for idx, (step_title, step_desc) in enumerate(flow_points):
    left = Inches(0.8 + idx * 2.4)
    top = Inches(2.2)
    width = Inches(2.2)
    height = Inches(4.2)
    
    # Draw block
    rect = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD_BG
    rect.line.color.rgb = BORDER_COLOR
    rect.line.width = Pt(1.5)
    
    tb = slide7.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_num = tf.paragraphs[0]
    p_num.text = f"STEP {idx+1}"
    p_num.font.name = "Segoe UI"
    p_num.font.size = Pt(12)
    p_num.font.bold = True
    p_num.font.color.rgb = ACCENT_ORANGE
    p_num.space_after = Pt(8)
    
    p_title = tf.add_paragraph()
    p_title.text = step_title
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(14)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_after = Pt(12)
    
    p_desc = tf.add_paragraph()
    p_desc.text = step_desc
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = TEXT_MUTED

# ----------------- SLIDE 8: Code Walkthrough & Snippets -----------------
slide8 = create_slide_with_bg()
add_slide_header(slide8, "Key Code Walkthrough & Server Snippets", "CODEBASE BASELINE")

# Left Column Code Card
left_code_title = "JWT Auth Handler (authController.js)"
left_code_text = """// Find user & Verify Password hash
const user = await prisma.user.findUnique({
    where: { email }
});
if (!user) {
    return res.status(401).json({
        message: 'Invalid credentials'
    });
}
const isValid = await bcrypt.compare(
    password, 
    user.passwordHash
);
if (!isValid) {
    return res.status(401).json({
        message: 'Invalid credentials'
    });
}
// Generate Token
const token = jwt.sign(
    { id: user.id, role: user.role },
    JWT_SECRET,
    { expiresIn: '30d' }
);"""

# Draw background code card left
rect_code_l = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
rect_code_l.fill.solid()
rect_code_l.fill.fore_color.rgb = CODE_BG
rect_code_l.line.color.rgb = BORDER_COLOR
rect_code_l.line.width = Pt(1.5)

tb_l = slide8.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.4))
tf_l = tb_l.text_frame
tf_l.word_wrap = True
tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0

p_c_title_l = tf_l.paragraphs[0]
p_c_title_l.text = left_code_title
p_c_title_l.font.name = "Segoe UI"
p_c_title_l.font.size = Pt(14)
p_c_title_l.font.bold = True
p_c_title_l.font.color.rgb = ACCENT_ORANGE
p_c_title_l.space_after = Pt(10)

p_code_l = tf_l.add_paragraph()
p_code_l.text = left_code_text
p_code_l.font.name = "Consolas"
p_code_l.font.size = Pt(9.5)
p_code_l.font.color.rgb = TEXT_WHITE

# Right Column Code Card
right_code_title = "Prisma Scheduler Query (schedulerService.js)"
right_code_text = """// Find overdue, uncompleted tasks
const overdueTasks = await prisma.task.findMany({
    where: {
        status: { not: 'COMPLETED' },
        dueDate: { lt: now }, // less than now
        employeeId: { not: null }
    },
    include: {
        employee: true,
        project: true
    }
});

// Send notification & email alert
for (const task of overdueTasks) {
    await prisma.email.create({
        data: {
            senderId: sender.id,
            receiverId: task.employeeId,
            subject: `Overdue Task: ${task.title}`
        }
    });
    // Triggers Nodemailer transport next...
}"""

# Draw background code card right
rect_code_r = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
rect_code_r.fill.solid()
rect_code_r.fill.fore_color.rgb = CODE_BG
rect_code_r.line.color.rgb = BORDER_COLOR
rect_code_r.line.width = Pt(1.5)

tb_r = slide8.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.4))
tf_r = tb_r.text_frame
tf_r.word_wrap = True
tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = 0

p_c_title_r = tf_r.paragraphs[0]
p_c_title_r.text = right_code_title
p_c_title_r.font.name = "Segoe UI"
p_c_title_r.font.size = Pt(14)
p_c_title_r.font.bold = True
p_c_title_r.font.color.rgb = ACCENT_ORANGE
p_c_title_r.space_after = Pt(10)

p_code_r = tf_r.add_paragraph()
p_code_r.text = right_code_text
p_code_r.font.name = "Consolas"
p_code_r.font.size = Pt(9.5)
p_code_r.font.color.rgb = TEXT_WHITE

# Save presentation
prs.save("application_overview.pptx")
print("Presentation compiled successfully as 'application_overview.pptx'")
